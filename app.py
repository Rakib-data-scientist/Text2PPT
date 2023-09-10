import os
import base64
import openai
import pptx
import streamlit as st
from dotenv import load_dotenv
from pptx.util import Inches, Pt

# Load environment variables
load_dotenv()
openai.api_key = os.getenv('sk-****')


# Define constants
TITLE_FONT_SIZE = Pt(30)
CONTENT_FONT_SIZE = Pt(16)
ENGINE = "text-davinci-003"

# Initialize Streamlit
st.title("Text2PPT Generation App")

def generate_slide_titles(topic):
    """Generates titles for slides based on the provided topic."""
    prompt = f"Generate 5 slide titles for the topic '{topic}'."
    response = openai.Completion.create(engine=ENGINE, prompt=prompt, max_tokens=200)
    return list(filter(bool, response.choices[0].text.strip().split("\n")))

def generate_slide_content(slide_title):
    """Generates content for a slide based on the provided title."""
    prompt = f"Generate content for the slide titled: '{slide_title}'."
    response = openai.Completion.create(engine=ENGINE, prompt=prompt, max_tokens=500)
    return response.choices[0].text.strip()

def create_presentation(topic, slide_titles_and_content):
    """Creates a PowerPoint presentation."""
    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = topic

    for title, content in slide_titles_and_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        slide.shapes.placeholders[1].text = content

        for paragraph in slide.shapes.placeholders[1].text_frame.paragraphs:
            paragraph.font.size = CONTENT_FONT_SIZE

    output_path = f"generated_ppt/{topic}_presentation.pptx"
    prs.save(output_path)
    return output_path

def get_download_link(file_path):
    """Generates a download link for the created PowerPoint file."""
    with open(file_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{os.path.basename(file_path)}">Download the PowerPoint Presentation</a>'

def main():
    topic = st.text_input("Enter the topic for your presentation:")
    if st.button("Generate Presentation") and topic:
        st.info("Generating presentation... Please wait.")
        
        slide_titles = generate_slide_titles(topic)
        slide_titles_and_content = [(title, generate_slide_content(title)) for title in slide_titles]
        
        output_path = create_presentation(topic, slide_titles_and_content)
        st.success("Presentation generated successfully!")
        st.markdown(get_download_link(output_path), unsafe_allow_html=True)

if __name__ == "__main__":
    main()
